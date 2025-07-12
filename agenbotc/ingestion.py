import uuid
import PyPDF2
import docx
from logging import Logger
import os
import hashlib
import json
from dotenv import load_dotenv
from typing import List, Dict, Any, Optional
import requests
from bs4 import BeautifulSoup
from langchain_text_splitters import RecursiveCharacterTextSplitter
from vector_store import add_to_vector_store, check_document_exists

# For backwards compatibility, keep the pptx import
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def calculate_content_hash(content):
    """Calculate MD5 hash of content string"""
    return hashlib.md5(content.encode('utf-8')).hexdigest()

def generate_doc_id_from_source(source, doc_type):
    """Generate a consistent document ID from source and type"""
    # Create a consistent ID based on source and type
    source_hash = hashlib.md5(f"{source}_{doc_type}".encode('utf-8')).hexdigest()
    return f"{doc_type}_{source_hash[:12]}"

def is_duplicate_document(source, doc_type):
    """Check if a document already exists in the vector store"""
    doc_id = generate_doc_id_from_source(source, doc_type)
    exists = check_document_exists(doc_id)
    return exists, doc_id if exists else None


def extract_text_from_pdf(pdf_path):
    """Extract text content from a PDF file"""
    text = ""
    try:
        # Check if file exists
        if not os.path.exists(pdf_path):
            raise ValueError(f"PDF file not found: {pdf_path}")
        
        # Check if file is readable
        if not os.access(pdf_path, os.R_OK):
            raise ValueError(f"PDF file is not readable: {pdf_path}")
        
        with open(pdf_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Check if PDF has pages
            if len(pdf_reader.pages) == 0:
                raise ValueError("PDF file contains no pages")
            
            for page in pdf_reader.pages:
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                except Exception as e:
                    print(f"Warning: Could not extract text from a page in {pdf_path}: {e}")
                    continue
        
        # Check if we extracted any text
        if not text.strip():
            raise ValueError("No text content could be extracted from the PDF")
            
        return text
        
    except Exception as e:
        raise ValueError(f"Error extracting text from PDF {pdf_path}: {str(e)}")


def extract_text_from_docx(file_path):
    """Extract text content from a Word file"""
    text = ""
    try:
        # Check if file exists
        if not os.path.exists(file_path):
            raise ValueError(f"DOCX file not found: {file_path}")
        
        # Check if file is readable
        if not os.access(file_path, os.R_OK):
            raise ValueError(f"DOCX file is not readable: {file_path}")
        
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():  # Only add non-empty paragraphs
                text += para.text + "\n"
        
        # Check if we extracted any text
        if not text.strip():
            raise ValueError("No text content could be extracted from the DOCX file")
            
        return text
        
    except Exception as e:
        raise ValueError(f"Error extracting text from DOCX {file_path}: {str(e)}")


def extract_text_from_website(url):
    """Extract text content from a website"""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        
        # Add timeout and error handling
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()  # Raises an HTTPError for bad responses
        
        # Check if response has content
        if not response.content:
            raise ValueError("Website returned empty content")
        
        soup = BeautifulSoup(response.content, "html.parser")
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.extract()
        
        # Get text content
        text = soup.get_text(separator="\n", strip=True)
        
        # Check if we extracted any meaningful text
        if not text or len(text.strip()) < 10:
            raise ValueError("No meaningful text content found on the website")
        
        return text
        
    except requests.exceptions.RequestException as e:
        raise ValueError(f"Failed to fetch website: {str(e)}")
    except Exception as e:
        raise ValueError(f"Error processing website content: {str(e)}")

def extract_text_from_ppt(ppt_path):
    """Extract text content from a PPT file"""
    if Presentation is None:
        raise ValueError("python-pptx library is not installed")
    
    text = ""
    try:
        # Check if file exists
        if not os.path.exists(ppt_path):
            raise ValueError(f"PPT file not found: {ppt_path}")
        
        # Check if file is readable
        if not os.access(ppt_path, os.R_OK):
            raise ValueError(f"PPT file is not readable: {ppt_path}")
        
        presentation = Presentation(ppt_path)
        
        # Check if presentation has slides
        if len(presentation.slides) == 0:
            raise ValueError("PPT file contains no slides")
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        
        # Check if we extracted any text
        if not text.strip():
            raise ValueError("No text content could be extracted from the PPT file")
            
        return text
        
    except Exception as e:
        raise ValueError(f"Error extracting text from PPT {ppt_path}: {str(e)}")

def chunk_text(text):
    """Split text into manageable chunks for processing"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    chunks = text_splitter.split_text(text)
    return chunks

def process_pdf(pdf_path):
    """Process a PDF file and add its content to the vector store with duplicate check"""
    try:
        # Check for duplicates using vector store
        is_duplicate, existing_doc_id = is_duplicate_document(pdf_path, "pdf")
        if is_duplicate:
            return {"doc_id": existing_doc_id, "message": "PDF already processed", "is_duplicate": True}
        
        # Generate consistent doc ID
        doc_id = generate_doc_id_from_source(pdf_path, "pdf")
        
        # Extract text with error handling
        try:
            text = extract_text_from_pdf(pdf_path)
        except Exception as e:
            return {"doc_id": None, "message": f"Failed to extract text from PDF: {str(e)}", "is_duplicate": False, "error": True}
        
        # Check if we extracted any meaningful text
        if not text or len(text.strip()) < 10:
            return {"doc_id": None, "message": "No meaningful text content found in PDF", "is_duplicate": False, "error": True}
        
        chunks = chunk_text(text)
        
        # Add chunks to vector store with metadata
        metadata = {
            "source": pdf_path,
            "type": "pdf",
            "doc_id": doc_id,
            "filename": os.path.basename(pdf_path)
        }
        add_to_vector_store(chunks, metadata)
        
        return {"doc_id": doc_id, "message": "PDF processed successfully", "is_duplicate": False}
        
    except Exception as e:
        return {"doc_id": None, "message": f"Unexpected error processing PDF: {str(e)}", "is_duplicate": False, "error": True}

def process_docx(docx_path):
    """Process a DOCX file and add its content to the vector store with duplicate check"""
    try:
        # Check for duplicates using vector store
        is_duplicate, existing_doc_id = is_duplicate_document(docx_path, "docx")
        if is_duplicate:
            return {"doc_id": existing_doc_id, "message": "DOCX already processed", "is_duplicate": True}
        
        # Generate consistent doc ID
        doc_id = generate_doc_id_from_source(docx_path, "docx")
        
        # Extract text with error handling
        try:
            text = extract_text_from_docx(docx_path)
        except Exception as e:
            return {"doc_id": None, "message": f"Failed to extract text from DOCX: {str(e)}", "is_duplicate": False, "error": True}
        
        # Check if we extracted any meaningful text
        if not text or len(text.strip()) < 10:
            return {"doc_id": None, "message": "No meaningful text content found in DOCX", "is_duplicate": False, "error": True}
        
        chunks = chunk_text(text)
        
        # Add chunks to vector store with metadata
        metadata = {
            "source": docx_path,
            "type": "docx",
            "doc_id": doc_id,
            "filename": os.path.basename(docx_path)
        }
        add_to_vector_store(chunks, metadata)
        
        print(f"Processed DOCX file: {docx_path}")
        return {"doc_id": doc_id, "message": "DOCX processed successfully", "is_duplicate": False}
        
    except Exception as e:
        return {"doc_id": None, "message": f"Unexpected error processing DOCX: {str(e)}", "is_duplicate": False, "error": True}

def process_website(url):
    """Process a website and add its content to the vector store with duplicate check"""
    try:
        # Check for duplicates using vector store
        is_duplicate, existing_doc_id = is_duplicate_document(url, "website")
        if is_duplicate:
            return {"doc_id": existing_doc_id, "message": "Website already processed", "is_duplicate": True}
        
        # Generate consistent doc ID
        doc_id = generate_doc_id_from_source(url, "website")
        
        # Extract text with error handling
        try:
            text = extract_text_from_website(url)
        except ValueError as e:
            return {"doc_id": None, "message": f"Failed to process website: {str(e)}", "is_duplicate": False, "error": True}
        
        chunks = chunk_text(text)
        
        # Add chunks to vector store with metadata
        metadata = {
            "source": url,
            "type": "website",
            "doc_id": doc_id
        }
        add_to_vector_store(chunks, metadata)
        
        return {"doc_id": doc_id, "message": "Website processed successfully", "is_duplicate": False}
        
    except Exception as e:
        return {"doc_id": None, "message": f"Unexpected error processing website: {str(e)}", "is_duplicate": False, "error": True}

def process_ppt(ppt_path):
    """Process a PPT file and add its content to the vector store with duplicate check"""
    try:
        # Check for PPT library availability
        if Presentation is None:
            return {"doc_id": None, "message": "PPT processing not available - python-pptx not installed", "is_duplicate": False, "error": True}
        
        # Check for duplicates using vector store
        is_duplicate, existing_doc_id = is_duplicate_document(ppt_path, "ppt")
        if is_duplicate:
            return {"doc_id": existing_doc_id, "message": "PPT already processed", "is_duplicate": True}
        
        # Generate consistent doc ID
        doc_id = generate_doc_id_from_source(ppt_path, "ppt")
        
        # Extract text with error handling
        try:
            text = extract_text_from_ppt(ppt_path)
        except Exception as e:
            return {"doc_id": None, "message": f"Failed to extract text from PPT: {str(e)}", "is_duplicate": False, "error": True}
        
        # Check if we extracted any meaningful text
        if not text or len(text.strip()) < 10:
            return {"doc_id": None, "message": "No meaningful text content found in PPT", "is_duplicate": False, "error": True}
        
        chunks = chunk_text(text)
        
        # Add chunks to vector store with metadata
        metadata = {
            "source": ppt_path,
            "type": "ppt",
            "doc_id": doc_id,
            "filename": os.path.basename(ppt_path)
        }
        add_to_vector_store(chunks, metadata)
        
        return {"doc_id": doc_id, "message": "PPT processed successfully", "is_duplicate": False}
        
    except Exception as e:
        return {"doc_id": None, "message": f"Unexpected error processing PPT: {str(e)}", "is_duplicate": False, "error": True}

